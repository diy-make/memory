from pptx import Presentation

try:
    prs = Presentation("x402 Builders Slides.pptx")
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(f"Shape: {shape.text[:50]}...")
except Exception as e:
    print(f"Error: {e}")
