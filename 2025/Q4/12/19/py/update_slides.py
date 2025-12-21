from pptx import Presentation
from pptx.util import Inches, Pt

def replace_text(slide, title_text, body_text):
    # Try to find title
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    
    # Try to find a body placeholder
    found_body = False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape == slide.shapes.title:
            continue
        
        # Assume the first non-title text frame is the body
        text_frame = shape.text_frame
        text_frame.clear()  # Clear existing text
        
        # Add paragraphs
        for line in body_text.split('\n'):
            p = text_frame.add_paragraph()
            p.text = line
            p.level = 0
            if line.startswith('  -'):
                p.level = 1
                p.text = line.replace('  -', '').strip()
            elif line.startswith('*'):
                p.text = line.replace('*', '').strip()

        found_body = True
        break
    
    if not found_body:
        # Create a text box if no body placeholder found
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = body_text

def add_slide(prs, title_text, body_text):
    # Use layout 1 (usually Title and Content)
    slide_layout = prs.slide_layouts[1] 
    slide = prs.slides.add_slide(slide_layout)
    replace_text(slide, title_text, body_text)

prs = Presentation("x402 Builders Slides.pptx")

# Content Data
slides_content = [
    (
        "OSO Token of Love & x402",
        "The Non-Pecuniary 'Joyfork' Protocol\n\nFrom Payment Required to Gratitude Accepted"
    ),
    (
        "A Brief History of HTTP 402",
        "x402 is not the first attempt to solve the 'Payment Required' puzzle.\n\n* WebMonetization.org & Interledger Protocol (ILP):\n  - Focus on 'streaming' payments (fractions of a cent).\n  - Passive browser integration via <link rel='monetization'>.\n* Grant For The Web:\n  - $100M fund to support open standards.\n* The Difference:\n  - Previous attempts focused on Creator Compensation.\n  - x402 focuses on Machine-to-Machine negotiation."
    ),
    (
        "The Webhook Nature (1/2)",
        "Generic Gateway vs. Protocol\n\n* Generic Webhook Gateway:\n  - A simple latch: Blocks access (401/403) until an event occurs.\n  - Problem: The client hits a wall and doesn't know *how* to unlock it.\n* The x402 Protocol:\n  - A standardized 'bridge'.\n  - It doesn't just block; it instructs.\n  - Returns 402 with WWW-Authenticate headers detailing exactly *who*, *what*, and *how much* to pay."
    ),
    (
        "The Webhook Nature (2/2)",
        "The 4-Turn Formalism (The 'Ping-Pong')\n\n* Turn 1 (Client): Request Resource (GET /image)\n* Turn 2 (Server): Challenge (402 Payment Required)\n  - Header: 'Send Love here'\n* Turn 3 (Client): Proof (Retry Request)\n  - Header: 'Here is the proof of Love'\n* Turn 4 (Server): Access (200 OK)\n  - Resource delivered."
    ),
    (
        "Liquity V2 'Joyfork'",
        "Friendly Forks vs. Happy Forks\n\n* Friendly Fork (e.g., Nerite):\n  - Licensed, serious financial protocol.\n  - Redemption = Stablecoin -> Collateral ($ to $).\n* Joyfork (OSO):\n  - A divergence in purpose.\n  - Redemption = Token of Love -> Nothing Monetary.\n  - We burn the token to create 'Flow' (Stream of Memories)."
    ),
    (
        "The OSO Vision",
        "Token of Love Streams\n\n* Closed-Loop Economy of Gratitude:\n  1. Creation: Audience gifts 'Love' (Tokens).\n  2. Usage: Robot 'spends' Love to work (x402).\n  3. Redemption: Tokens are burned to create a permanent record.\n* Result:\n  - The robot runs on vibes.\n  - A protocol for the AI Agent economy."
    )
]

# Overwrite existing slides
for i, (title, body) in enumerate(slides_content):
    if i < len(prs.slides):
        replace_text(prs.slides[i], title, body)
    else:
        add_slide(prs, title, body)

# Save
prs.save("x402 Builders Slides Updated.pptx")
print("Done.")
